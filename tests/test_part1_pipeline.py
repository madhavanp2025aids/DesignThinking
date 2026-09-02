"""
Part 1 Tests: Multi-Format Ingestion & Specification Extraction Pipeline
Verifies format parsers, extraction accuracy, verification grounding,
no-hallucination / not-available flags, and API endpoints.
"""

from fastapi.testclient import TestClient
from backend.main import app
from execution.parse_csv import parse_csv
from execution.parse_pptx import parse_pptx
from execution.spec_verifier import SpecVerifier
from execution.spec_extractor import extract_specs_from_document


def test_csv_parser(tmp_path):
    csv_file = tmp_path / "specs.csv"
    csv_file.write_text("Parameter,Value,Unit\nOuter Diameter,120,mm\nBore,80,mm\nLength,450,mm\nMaterial,AISI 316L,\n", encoding="utf-8")
    
    parsed = parse_csv(str(csv_file), "specs.csv")
    assert "Outer Diameter" in parsed["raw_text"]
    assert len(parsed["tables"]) == 1
    assert parsed["tables"][0]["headers"] == ["Parameter", "Value", "Unit"]
    assert len(parsed["tables"][0]["rows"]) == 4


def test_pptx_parser(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    pptx_file = tmp_path / "machine_specs.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = txBox.text_frame
    tf.text = "Shaft Diameter: 50 mm\nTotal Length: 300 mm\nMaterial: 42CrMo4"

    prs.save(str(pptx_file))

    parsed = parse_pptx(str(pptx_file), "machine_specs.pptx")
    assert len(parsed["slides"]) == 1
    assert "Shaft Diameter: 50 mm" in parsed["raw_text"]


def test_spec_verifier_grounding():
    source_text = "HYDAC Hydraulic Cylinder specs:\nOuter Diameter: 100 mm\nWorking Pressure: 250 bar\nMaterial: 42CrMo4"
    
    # Grounded values must pass
    valid, note = SpecVerifier.verify_field(source_text, "100 mm", "100", "Outer Diameter: 100 mm")
    assert valid is True

    # Ungrounded values must be rejected
    invalid, note = SpecVerifier.verify_field(source_text, "999 mm", "999", "Random guessed text")
    assert invalid is False


def test_spec_extractor_and_not_available_rule():
    parsed_mock = {
        "raw_text": "Bore Diameter: 80 mm\nRod Diameter: 45 mm\nStroke: 500 mm\nMaterial: Steel",
        "tables": [],
        "text_blocks": [
            {"text": "Bore Diameter: 80 mm\nRod Diameter: 45 mm\nStroke: 500 mm\nMaterial: Steel", "source_location": "Page 1"}
        ]
    }

    fields = extract_specs_from_document(
        parsed_data=parsed_mock,
        document_id="doc-123",
        part_id="part-456",
        filename="cylinder_spec.pdf",
        required_field_names=["bore_diameter", "rod_diameter", "stroke", "flange_diameter", "pitch"]
    )

    field_map = {f["field_name"]: f for f in fields}

    # Present fields must be available with full source trace
    assert field_map["bore_diameter"]["is_available"] == 1
    assert field_map["bore_diameter"]["normalized_value"] == "80"
    assert field_map["bore_diameter"]["unit"] == "mm"
    assert field_map["bore_diameter"]["source_location"] == "Page 1"

    # Missing fields must be explicitly marked is_available: 0 and "Not available in uploaded document"
    assert field_map["flange_diameter"]["is_available"] == 0
    assert field_map["flange_diameter"]["not_available_reason"] == "Not available in uploaded document"
    assert field_map["flange_diameter"]["raw_value"] is None


def test_api_spec_ingestion_and_extraction_flow(auth_headers):
    client = TestClient(app)
    headers = {"Authorization": auth_headers["Authorization"]}

    # 1. Create a Part
    res = client.post("/api/specs/parts", json={"name": "Hydraulic Ram CY-500"}, headers=headers)
    assert res.status_code == 201
    part_id = res.json()["id"]

    # 2. Upload CSV spec document for the part
    csv_content = b"Specification,Value,Unit\nOuter Diameter,140,mm\nInner Diameter,100,mm\nLength,600,mm\nMaterial,Stainless Steel 316,\n"
    res_upload = client.post(
        f"/api/specs/parts/{part_id}/documents",
        files=[("files", ("specs.csv", csv_content, "text/csv"))],
        data={"auto_extract": "true"},
        headers=headers
    )
    assert res_upload.status_code == 201
    assert len(res_upload.json()["documents"]) == 1

    # 3. Retrieve extracted fields
    res_fields = client.get(f"/api/specs/parts/{part_id}/fields", headers=headers)
    assert res_fields.status_code == 200
    fields = res_fields.json()
    assert len(fields) > 0

    field_by_name = {f["field_name"]: f for f in fields}
    assert "outer_diameter" in field_by_name
    assert field_by_name["outer_diameter"]["is_available"] is True
    assert field_by_name["outer_diameter"]["normalized_value"] == "140"
    assert field_by_name["outer_diameter"]["unit"] == "mm"

    # 4. Check Status
    res_status = client.get(f"/api/specs/parts/{part_id}/status", headers=headers)
    assert res_status.status_code == 200
    status_data = res_status.json()
    assert status_data["total_documents"] == 1
    assert status_data["available_fields"] >= 3

    # 5. User correction
    target_field = field_by_name["outer_diameter"]
    res_corr = client.put(
        f"/api/specs/fields/{target_field['id']}",
        json={"correction": "145", "unit": "mm"},
        headers=headers
    )
    assert res_corr.status_code == 200
    assert res_corr.json()["user_correction"] == "145"
    assert "140" in res_corr.json()["raw_value"]
