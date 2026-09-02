"""
Spec-to-3D Generator — PPTX / PPT Parser
Extracts text from slides, text shapes, embedded tables, and speaker notes.
Preserves slide number and element index for full source traceability.
"""

from typing import Optional


def parse_pptx(file_path: str, filename: str) -> dict:
    """
    Parse a PPTX file, extracting text boxes, tables, and speaker notes per slide.

    Returns:
        {
            "raw_text": str,
            "tables": [
                {
                    "table_index": int,
                    "source_location": str,
                    "headers": [str],
                    "rows": [[str]]
                }
            ],
            "slides": [
                {
                    "slide_num": int,
                    "text": str,
                    "tables": [...],
                    "notes": str
                }
            ]
        }
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx is required for PPTX parsing. Install: pip install python-pptx")

    prs = Presentation(file_path)

    result = {
        "raw_text": "",
        "tables": [],
        "slides": [],
    }

    table_index = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text_chunks = []
        slide_tables = []
        notes_text = ""

        # Extract shapes & tables in slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        slide_text_chunks.append(line)

            elif shape.has_table:
                table = shape.table
                rows_data = []
                for row in table.rows:
                    row_cells = [cell.text.strip() if cell.text else "" for cell in row.cells]
                    if any(cell != "" for cell in row_cells):
                        rows_data.append(row_cells)

                if len(rows_data) >= 1:
                    headers = rows_data[0]
                    data_rows = rows_data[1:] if len(rows_data) > 1 else [rows_data[0]]

                    table_entry = {
                        "table_index": table_index,
                        "source_location": f"Slide {slide_idx}, Table {table_index + 1}",
                        "headers": headers,
                        "rows": data_rows,
                    }
                    slide_tables.append(table_entry)
                    result["tables"].append(table_entry)
                    table_index += 1

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        combined_slide_text = "\n".join(slide_text_chunks)
        if notes_text:
            combined_slide_text += f"\n[Speaker Notes]: {notes_text}"

        slide_entry = {
            "slide_num": slide_idx,
            "text": combined_slide_text,
            "tables": slide_tables,
            "notes": notes_text,
        }

        result["slides"].append(slide_entry)
        result["raw_text"] += f"\n--- Slide {slide_idx} ---\n{combined_slide_text}\n"

    return result
